from pptx import Presentation
from .privileges import has_permission

def slides(prompt: str, user: str) -> str:
    """
    Generates a PowerPoint presentation based on the prompt.
    """
    if not has_permission(user, "create_slides"):
        return "You do not have permission to create slides."

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Presentation Title"
    subtitle.text = prompt

    # Save the presentation to a temporary file
    filepath = "/tmp/presentation.pptx"
    prs.save(filepath)

    return f"Presentation created at {filepath}"

def download_slide(user: str, slide_id: str) -> str:
    """
    Simulates downloading a slide if the user has the required privileges.
    """
    if not has_permission(user, "download_slides"):
        return "You do not have permission to download slides."

    # In a real application, you would fetch the slide from storage.
    return f"Downloading slide {slide_id}..."

def present_slide(user: str, slide_id: str) -> str:
    """
    Simulates presenting a slide if the user has the required privileges.
    """
    if not has_permission(user, "present_slides"):
        return "You do not have permission to present slides."

    # In a real application, you would display the slide in a presenter view.
    return f"Presenting slide {slide_id}..."
